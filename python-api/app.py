from flask import Flask, request, jsonify
from flask_cors import CORS
import requests
from bs4 import BeautifulSoup
import PyPDF2
import docx
import pptx

app = Flask(__name__)
CORS(app)

@app.route('/summarize', methods=['POST'])
def summarize():
    url = request.json.get('url')
    response = requests.get(url)
    soup = BeautifulSoup(response.content, 'html.parser')
    text = ' '.join([p.text for p in soup.find_all('p')])
    summary = summarize_text(text)
    return jsonify({'summary': summary})

@app.route('/upload', methods=['POST'])
def upload():
    file = request.files['file']
    if file.filename.endswith('.pdf'):
        reader = PyPDF2.PdfFileReader(file)
        text = ' '.join([reader.getPage(i).extract_text() for i in range(reader.numPages)])
    elif file.filename.endswith('.docx'):
        doc = docx.Document(file)
        text = ' '.join([para.text for para in doc.paragraphs])
    elif file.filename.endswith('.pptx'):
        presentation = pptx.Presentation(file)
        text = ' '.join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, 'text')])
    else:
        return jsonify({'error': 'Unsupported file type'}), 400
    summary = summarize_text(text)
    return jsonify({'summary': summary})

def summarize_text(text):
    # Implement your text summarization logic here
    return text[:200]  # Placeholder: return the first 200 characters

if __name__ == '__main__':
    app.run(debug=True)
